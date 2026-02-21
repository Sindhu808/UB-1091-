from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation(filename="HELIX_Project_Documentation.pptx"):
    prs = Presentation()

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "HELIX: Hybrid Renewable Energy VPP"
    subtitle.text = "Intelligent Virtual Power Plant Orchestration\nProject Documentation"

    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "I. The Problem Statement & Solution (1/2)"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "The Problem:"
    
    p = tf.add_paragraph()
    p.text = "• Renewable Energy Curtailment: Excess solar and wind energy is often wasted or curtailed when grid demand is low."
    
    p = tf.add_paragraph()
    p.text = "• Fossil Fuel Dependency: Grids are forced to rely on dirty backup generation (like coal or diesel) when weather restricts renewable output."
    
    p = tf.add_paragraph()
    p.text = "• Unorchestrated Grids: A lack of real-time intelligence at campus and regional boundaries leads to highly inefficient energy distribution and utilization."

    # Slide 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "I. The Problem Statement & Solution (2/2)"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "The Solution - HELIX:"
    
    p = tf.add_paragraph()
    p.text = "• Intelligent Virtual Power Plant (VPP): Unified orchestration of separate energy sources (Solar, Wind, Hydro, and Biomass) into a single smart micro-grid."
    
    p = tf.add_paragraph()
    p.text = "• Real-Time Arbitrage: Deep integration with battery storage to cache surplus energy instantly and supply deficits autonomously."
    
    p = tf.add_paragraph()
    p.text = "• Predictive AI-Driven Orchestration: Uses forward-looking 48-hour weather forecasts to proactively dictate charge & discharge cycles based on expected generation."

    # Slide 4: Tools and Technologies Used
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "II. Tools and Technologies Used"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Frontend (User Interface):"
    
    p = tf.add_paragraph()
    p.text = "• React.js & Vite: Powering a high-performance 60FPS dynamic dashboard UI."
    
    p = tf.add_paragraph()
    p.text = "• Vanilla CSS: Implemented with modern CSS Variables for a dynamic, sleek design system."
    
    p = tf.add_paragraph()
    p.text = "\nBackend (Core Engine & API):"
    
    p = tf.add_paragraph()
    p.text = "• Python with FastAPI: For lightning-fast REST Endpoints and concurrency."
    
    p = tf.add_paragraph()
    p.text = "• WebSockets: For real-time telemetry streaming straight to the frontend."
    
    p = tf.add_paragraph()
    p.text = "• SQLAlchemy & SQLite: For lightweight, real-time application state persistence."

    # Slide 5: Implementation Details
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "III. Implementation Details"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Technical Milestones:"
    
    p = tf.add_paragraph()
    p.text = "• Live Data Ingestion Pipeline: Processes live dataset feeds mimicking continuous 100ms IoT sensor polling."
    
    p = tf.add_paragraph()
    p.text = "• Optimization Engine: Autonomously adjusts battery charge/discharge rates using Time-of-Use pricing logic to dodge expensive grid tariffs during peak hours."
    
    p = tf.add_paragraph()
    p.text = "• Forecasting Integration: Consumes Open-Meteo forecasts to predict structural deficits up to 48 hours in advance."
    
    p = tf.add_paragraph()
    p.text = "• Actionable Recommendation Feed: Auto-generates real-time, human-readable insights for dashboard operators (e.g., 'Warning: Shifting HVAC load to Solar peak')."

    # Slide 6: Screenshots & Links
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "IV. Screenshots & Repository Links"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "[Please insert screenshots of the HELIX Dashboard here]"
    
    p = tf.add_paragraph()
    p.text = "\nGitHub Repository:"
    
    p = tf.add_paragraph()
    p.text = "• https://github.com/[your-username]/HELIX-VPP"
    
    p = tf.add_paragraph()
    p.text = "\nLive Hosted Demo:"
    
    p = tf.add_paragraph()
    p.text = "• Frontend Hosted on [Vercel/Netlify]: https://helix-[insert-app].vercel.app"
    
    p = tf.add_paragraph()
    p.text = "• Backend Hosted on [Render/Railway]"

    prs.save(filename)
    print(f"Presentation saved successfully as {filename}")

if __name__ == '__main__':
    create_presentation()
